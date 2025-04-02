# -*- coding: utf-8 -*-
"""Personal Chatbot

Automatically generated by Colab.

Original file is located at
    https://colab.research.google.com/drive/1KXG-RGbtlhNt2dWsy7hxtI0k9ZM2ZdQX
"""

pip install pdfplumber python-docx python-pptx faiss-cpu sentence-transformers openai gradio

import os
import pdfplumber
from docx import Document
from pptx import Presentation
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import gradio as gr
from openai import OpenAI

# Set up paths and constants
RESUME_FOLDER = "xx" # add your grive path

# Global counters to track file types processed
pdf_count = 0
docx_count = 0
pptx_count = 0

# Function to check if a file is a valid PDF
def is_valid_pdf(pdf_path):
    try:
        with open(pdf_path, "rb") as f:
            return f.read(4) == b"%PDF"
    except:
        return False

# Function to extract text from PDFs, ignoring corrupt ones
def extract_text_from_pdf(pdf_path):
    global pdf_count
    pdf_count += 1  # Increment PDF count
    print(f"PDFs processed: {pdf_count}, DOCX: {docx_count}, PPTX: {pptx_count}", flush=True)
    text = ""
    try:
        if is_valid_pdf(pdf_path):
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
        else:
            print(f"Skipping invalid PDF: {pdf_path}")
            return ""
    except Exception as e:
        print(f"Error reading {pdf_path}: {e}")
        return ""  # Return empty string if there's an error
    return text

# Function to check if a DOCX file is valid
def is_valid_docx(docx_path):
    try:
        Document(docx_path)  # Try opening the document
        return True
    except:
        return False

# Function to extract text from DOCX files, ignoring corrupt ones
def extract_text_from_docx(docx_path):
    global docx_count
    docx_count += 1  # Increment DOCX count
    print(f"PDFs processed: {pdf_count}, DOCX: {docx_count}, PPTX: {pptx_count}", flush=True)
    try:
        if is_valid_docx(docx_path):
            doc = Document(docx_path)
            return "\n".join([para.text for para in doc.paragraphs])
        else:
            print(f"Skipping invalid DOCX: {docx_path}")
            return ""
    except Exception as e:
        print(f"Error reading {docx_path}: {e}")
        return ""

# Function to extract text from PPTX files with error handling (already in place)
def extract_text_from_pptx(pptx_path):
    global pptx_count
    pptx_count += 1  # Increment PPTX count
    print(f"PDFs processed: {pdf_count}, DOCX: {docx_count}, PPTX: {pptx_count}", flush=True)
    try:
        text = ""
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading {pptx_path}: {e}")
        return ""  # Return empty string if there's an error

# Function to add filename to the extracted text
def add_filename_to_text(file_path, extracted_text):
    filename = os.path.basename(file_path)
    return f"Filename: {filename}\n{extracted_text}"

# Load and process files from the folder
def load_files_from_folder(folder_path):
    all_files = [os.path.join(root, file) for root, _, files in os.walk(folder_path) for file in files if file.endswith((".pdf", ".docx", ".pptx")) and not file.startswith("~$")]
    all_texts = []
    for file in all_files:
        if file.endswith(".pdf"):
            extracted_text = extract_text_from_pdf(file)
        elif file.endswith(".docx"):
            extracted_text = extract_text_from_docx(file)
        elif file.endswith(".pptx"):
            extracted_text = extract_text_from_pptx(file)

        if extracted_text:  # Only add if there's extracted text
            all_texts.append(add_filename_to_text(file, extracted_text))
    return all_texts, all_files

# Load all the text from files
file_texts, file_names = load_files_from_folder(RESUME_FOLDER)

import json

with open("/content/file_data.json", "w") as f:
    json.dump({"file_texts": file_texts, "file_names": file_names}, f)

from google.colab import files
files.download("/content/file_data.json")

!pip install huggingface_hub

!huggingface-cli login

import json
import torch
import gradio as gr
from openai import OpenAI
from sentence_transformers import SentenceTransformer
import faiss
import os

# Load extracted file data from JSON
with open("file_data.json", "r", encoding="utf-8") as f:
    file_data = json.load(f)

file_texts = file_data["file_texts"]
file_names = file_data["file_names"]

# Load embedding model with GPU support
device = "cuda" if torch.cuda.is_available() else "cpu"
embedding_model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2", device=device)

# Encode text embeddings
embeddings = embedding_model.encode(file_texts, convert_to_numpy=True)

# Initialize FAISS index (GPU if available, otherwise CPU)
dimension = embeddings.shape[1]
cpu_index = faiss.IndexFlatL2(dimension)  # CPU index

try:
    res = faiss.StandardGpuResources()  # Try initializing GPU resources
    gpu_index = faiss.index_cpu_to_gpu(res, 0, cpu_index)  # Move index to GPU
    gpu_index.add(embeddings)  # Add embeddings to FAISS index
    index = gpu_index  # Use GPU index
    print("Using FAISS with GPU acceleration.")
except AttributeError:
    print("FAISS GPU not available. Using CPU.")
    cpu_index.add(embeddings)  # Use CPU FAISS
    index = cpu_index  # Use CPU index

# Function to retrieve relevant documents
def retrieve_documents(query):
    query_embedding = embedding_model.encode([query], convert_to_numpy=True)
    D, I = index.search(query_embedding, k=5)  # Search top 5 matches
    return [(file_texts[idx], os.path.basename(file_names[idx])) for idx in I[0]]  # Extract only filename

# OpenAI client setup
client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key="xx" #enter your api key
)

# Function for refining responses with sources included
def refine_response(user_input):
    retrieved_docs = retrieve_documents(user_input)
    if not retrieved_docs:
        return "No relevant documents found."

    context_text = "\n\n".join([doc[0] for doc in retrieved_docs])
    source_files = ", ".join([doc[1] for doc in retrieved_docs])  # Collect sources

    # Format response with sources
    response_text = f"**Answer Based on Retrieved Documents:**\n\n{context_text}\n\n**Sources:** {source_files}"

    # Call LLM for enhanced response
    completion = client.chat.completions.create(
        model="deepseek/deepseek-r1-zero:free",
        messages=[
            {"role": "system", "content": "You are an AI assistant. Answer the user's query based on the provided context and cite the relevant document names."},
            {"role": "user", "content": f"Query: {user_input}\n\nContext:\n{context_text}\n\nCite sources: {source_files}"}
        ]
    )

    llm_response = completion.choices[0].message.content
    return f"{llm_response}\n\n**Sources:** {source_files}"


# Define the chatbot function
def chatbot(user_input, history):
    return refine_response(user_input)

# Create and launch the Gradio Chat Interface
demo = gr.ChatInterface(chatbot)
demo.launch(share=True, debug=True)